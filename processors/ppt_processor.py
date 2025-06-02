from pptx import Presentation
from typing import Dict, Any, List
from .base_processor import BaseProcessor

class PPTProcessor(BaseProcessor):
    """Process PowerPoint files"""
    
    def process(self, filepath: str) -> Dict[str, Any]:
        """Extract text and structure from PowerPoint presentations"""
        result = {
            'text': '',
            'slides': [],
            'metadata': {}
        }
        
        try:
            prs = Presentation(filepath)
            
            # Extract metadata
            core_props = prs.core_properties
            result['metadata'] = {
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'created': str(core_props.created) if core_props.created else '',
                'modified': str(core_props.modified) if core_props.modified else '',
                'slide_count': len(prs.slides)
            }
            
            text_parts = [f"# Presentation\n\n**Slides:** {len(prs.slides)}\n\n"]
            
            # Process each slide
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_data = {
                    'slide_number': slide_idx,
                    'title': '',
                    'content': [],
                    'notes': ''
                }
                
                slide_text = []
                slide_text.append(f"## Slide {slide_idx}\n\n")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        
                        # Try to identify title
                        if not slide_data['title'] and (
                            shape.shape_type == 14 or  # Title placeholder
                            len(text.split('\n')[0]) < 100  # Short first line likely a title
                        ):
                            slide_data['title'] = text.split('\n')[0]
                            slide_text.append(f"### {slide_data['title']}\n\n")
                            
                            # Add remaining content if any
                            remaining_text = '\n'.join(text.split('\n')[1:]).strip()
                            if remaining_text:
                                slide_data['content'].append(remaining_text)
                                slide_text.append(f"{remaining_text}\n\n")
                        else:
                            slide_data['content'].append(text)
                            slide_text.append(f"{text}\n\n")
                
                # Extract notes
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    for shape in notes_slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            slide_data['notes'] = shape.text.strip()
                            slide_text.append(f"**Speaker Notes:** {slide_data['notes']}\n\n")
                            break
                
                result['slides'].append(slide_data)
                text_parts.extend(slide_text)
            
            result['text'] = '\n'.join(text_parts)
            
        except Exception as e:
            result['error'] = str(e)
            result['text'] = f"Error processing PowerPoint: {str(e)}"
        
        return result
    
    def to_markdown(self, content: Dict[str, Any]) -> str:
        """Convert PowerPoint content to markdown"""
        if isinstance(content, str):
            return content
        
        if content.get('text'):
            return content['text']
        
        # Fallback conversion
        markdown = "# Presentation\n\n"
        
        # Add metadata
        if content.get('metadata'):
            metadata = content['metadata']
            markdown += "## Information\n\n"
            for key, value in metadata.items():
                if value:
                    markdown += f"**{key.title()}:** {value}\n\n"
        
        # Add slides
        if content.get('slides'):
            markdown += "## Slides\n\n"
            for slide in content['slides']:
                markdown += f"### Slide {slide['slide_number']}"
                
                if slide.get('title'):
                    markdown += f": {slide['title']}\n\n"
                else:
                    markdown += "\n\n"
                
                if slide.get('content'):
                    for content_item in slide['content']:
                        markdown += f"{content_item}\n\n"
                
                if slide.get('notes'):
                    markdown += f"**Speaker Notes:** {slide['notes']}\n\n"
                
                markdown += "---\n\n"
        
        return markdown
